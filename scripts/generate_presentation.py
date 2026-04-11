#!/usr/bin/env python3
"""SQiP 2026 論文説明用パワーポイント生成"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = "/Users/suetaketakaya/SoftwareQualitySymposium"
FIG = os.path.join(BASE, "drafts", "figures")
OUT = os.path.join(BASE, "report", "presentation_2026.pptx")

# Colors
BLACK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
ACCENT = RGBColor(0x1A, 0x56, 0x8E)  # Deep blue
ACCENT2 = RGBColor(0xC0, 0x39, 0x2B)  # Red for emphasis
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SECTION = RGBColor(0x16, 0x3C, 0x6E)

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Hiragino Sans"):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet(tf, text, font_size=16, bold=False, color=BLACK, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Hiragino Sans"
    p.level = level
    p.space_after = Pt(4)
    return p

def add_image(slide, path, left, top, width):
    slide.shapes.add_picture(path, Cm(left), Cm(top), Cm(width))

def add_rect(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Hiragino Sans"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)
    return shape


# =====================================================================
# Slide 1: Title
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, BG_DARK)

add_text_box(slide, 2, 2, 30, 2,
    "SQiP 2026 経験発表", 20, color=RGBColor(0xAA,0xAA,0xAA))
add_text_box(slide, 2, 4.5, 30, 4,
    "静的解析に基づくテスト要求モデルの自動生成と\nそれを用いたLLM単体テスト生成の実証評価",
    28, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_text_box(slide, 2, 10, 30, 2,
    "品質管理・テスト技術の観点", 18, color=RGBColor(0x88,0xBB,0xDD))
add_text_box(slide, 2, 13, 30, 2,
    "ホワイトボックステスト / テスト要求モデル / LLM / 静的解析 / 単体テスト生成 / トレーサビリティ",
    14, color=GRAY)

# Key numbers
for i, (num, label) in enumerate([("99件", "テスト要求"), ("145件", "生成テスト"),
                                    ("91.7%", "初回PASS率"), ("100%", "修正後PASS率")]):
    x = 2 + i * 7.5
    add_rect(slide, x, 15, 6.5, 2.5, ACCENT, f"{num}\n{label}", 16)


# =====================================================================
# Slide 2: 背景と問題
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.8, 31, 1.5, "1. 背景: LLM活用開発における単体テストの課題", 24, bold=True, color=ACCENT)

# Left: Current situation
add_rect(slide, 1.5, 3, 14.5, 3.5, RGBColor(0xE8,0xF0,0xF8), "", 12)
add_text_box(slide, 2, 3.2, 13.5, 0.8, "LLM活用開発の現状", 16, bold=True, color=ACCENT)
tf = add_text_box(slide, 2, 4.2, 13.5, 2, "", 14, color=BLACK)
add_bullet(tf, "荒い要求 → LLMがコード生成 → 基本機能の動作確認", 13, color=GRAY)
add_bullet(tf, "評価が「動くかどうか」に偏りがち", 13, color=GRAY)
add_bullet(tf, "単体テスト・コンポーネントテスト層が手薄", 13, bold=True, color=ACCENT2)

# Right: What's needed
add_rect(slide, 17.5, 3, 14.5, 3.5, RGBColor(0xFD,0xF0,0xF0), "", 12)
add_text_box(slide, 18, 3.2, 13.5, 0.8, "実リリースに必要なこと", 16, bold=True, color=ACCENT2)
tf = add_text_box(slide, 18, 4.2, 13.5, 2, "", 14, color=BLACK)
add_bullet(tf, "各関数が仕様どおりに動作することの検証", 13, color=GRAY)
add_bullet(tf, "分岐条件・境界値の体系的なテスト", 13, color=GRAY)
add_bullet(tf, "テスト設計の根拠と網羅性の管理", 13, bold=True, color=ACCENT2)

# P1/P2/P3
add_text_box(slide, 1.5, 7.5, 31, 1, "既存LLMテスト生成手法の3つの問題", 18, bold=True, color=BLACK)

problems = [
    ("P1", "テスト設計根拠の不在", "生成テストが「何の分岐を」「どの同値クラスを」対象としているか不明。\n第三者がテスト設計の妥当性を検証できない。"),
    ("P2", "網羅性の事前把握が不可能", "テスト生成「前」に網羅性を設計・管理する手段がない。\n重複テストばかりという状況を事前に検知できない。"),
    ("P3", "テスト失敗時の原因特定コスト", "テスト失敗時にテストコードとソースコードを逐一突合する必要。\n145件のテスト失敗原因の特定に長時間を要する。"),
]
for i, (pid, title, desc) in enumerate(problems):
    x = 1.5 + i * 10.5
    add_rect(slide, x, 9, 9.5, 1.2, ACCENT2 if i==0 else ACCENT, f"{pid}: {title}", 13)
    add_text_box(slide, x + 0.3, 10.5, 9, 2.5, desc, 11, color=GRAY)

# Bottom: prior work note
add_text_box(slide, 1.5, 13.5, 31, 2,
    "先行研究 (CodaMOSA, ChatUniTest, TestPilot, CoverUp等) はカバレッジ向上に注力。\n"
    "115本のサーベイ[7]で89%が「直接テスト生成」方式 → テスト設計情報が外部化されない",
    12, color=GRAY)


# =====================================================================
# Slide 3: 提案手法
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.8, 31, 1.5, "2. 提案手法: テスト要求モデル (TRM) による3段階パイプライン", 22, bold=True, color=ACCENT)

add_image(slide, os.path.join(FIG, "fig1-overview.png"), 1.5, 2.5, 30)

# Key point
add_rect(slide, 1.5, 13.5, 30, 3.5, RGBColor(0xE8,0xF0,0xF8), "", 12)
add_text_box(slide, 2, 13.7, 29, 0.8, "核心: TRM（テスト要求モデル）= 構造化された中間成果物", 16, bold=True, color=ACCENT)
tf = add_text_box(slide, 2, 14.7, 29, 2, "", 13, color=BLACK)
add_bullet(tf, "分岐網羅(BR)・同値クラス(EC)・境界値(BV)・エラーパス(ER)・依存切替(DP) の5種別", 13)
add_bullet(tf, "各要求にID・説明・優先度・ソースコード参照を付与 → YAML形式で出力", 13)
add_bullet(tf, "P1解決: テスト設計根拠が追跡可能 / P2解決: 生成「前」に網羅性管理 / P3解決: FAIL時にIDで即時診断", 13, bold=True, color=ACCENT)


# =====================================================================
# Slide 4: TRMの具体例
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.8, 31, 1.5, "TRMの具体例: ParseVersion関数", 22, bold=True, color=ACCENT)

# YAML example
add_rect(slide, 1.5, 2.5, 15, 7, RGBColor(0xF5,0xF5,0xF5), "", 12)
add_text_box(slide, 2, 2.7, 14, 0.8, "Step 2 出力: TRM (YAML)", 14, bold=True, color=ACCENT)
yaml_text = (
    '- id: "BR-11"\n'
    '  type: branch_coverage\n'
    '  description: "alpha修飾子の\n'
    '    分岐を通過する"\n'
    '  priority: high\n'
    '  source_ref: "BC-02-10:\n'
    '    strncmp(p,\\"alpha\\",5)==0"\n'
    '  expected_behavior:\n'
    '    "nShift = -0x60 を設定"'
)
add_text_box(slide, 2.2, 3.8, 13.5, 6, yaml_text, 11, color=BLACK, font_name="Consolas")

# Test code example
add_rect(slide, 17.5, 2.5, 15, 7, RGBColor(0xF5,0xF5,0xF5), "", 12)
add_text_box(slide, 18, 2.7, 14, 0.8, "Step 3 出力: テストコード (C++)", 14, bold=True, color=ACCENT)
code_text = (
    'TEST(ParseVersion,\n'
    '     AlphaModifier) {\n'
    '  // BR-11: alpha修飾子の\n'
    '  //   分岐を通過する\n'
    '  UINT32 val =\n'
    '    ParseVersion(\n'
    '      L"2.4.1alpha");\n'
    '  EXPECT_EQ(val,\n'
    '    0x82848120);\n'
    '}'
)
add_text_box(slide, 18.2, 3.8, 13.5, 6, code_text, 11, color=BLACK, font_name="Consolas")

# Arrow
add_text_box(slide, 15.5, 5.5, 2, 1, "  =>", 28, bold=True, color=ACCENT)

# Traceability point
add_rect(slide, 1.5, 10.5, 30.5, 2.5, ACCENT, "", 12)
add_text_box(slide, 2, 10.7, 29, 0.7, "TRM IDによるトレーサビリティ", 16, bold=True, color=WHITE)
tf = add_text_box(slide, 2, 11.7, 29, 1.2, "", 13, color=RGBColor(0xCC,0xDD,0xEE))
add_bullet(tf, "テスト要求 BR-11 → テストケース ParseVersion.AlphaModifier → 正方向・逆方向の追跡が可能", 13, color=RGBColor(0xCC,0xDD,0xEE))
add_bullet(tf, "テスト失敗時: BR-11 の説明とソースコード参照から原因を即座に特定", 13, color=WHITE)

# 5 types
add_text_box(slide, 1.5, 14, 31, 0.8, "5種別の体系", 16, bold=True, color=BLACK)
types_data = [("BR", "分岐網羅\n55件"), ("EC", "同値クラス\n27件"), ("BV", "境界値\n11件"), ("ER", "エラーパス\n3件"), ("DP", "依存切替\n3件")]
for i, (sym, label) in enumerate(types_data):
    x = 1.5 + i * 6.2
    shade = max(0x30, 0x88 - i * 0x15)
    c = RGBColor(shade, shade, shade + 0x20)
    add_rect(slide, x, 15, 5.5, 2.5, c, f"{sym}\n{label}", 13)


# =====================================================================
# Slide 5: 実験条件
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.8, 31, 1.5, "3. 実験条件", 24, bold=True, color=ACCENT)

# Left: Target
add_text_box(slide, 1.5, 2.5, 15, 1, "対象プロジェクト", 16, bold=True, color=BLACK)
tf = add_text_box(slide, 1.5, 3.5, 15, 6, "", 13, color=GRAY)
add_bullet(tf, "sakura-editor/sakura (C++, OSS)", 13, color=BLACK)
add_bullet(tf, "Windows向け日本語テキストエディタ", 13)
add_bullet(tf, "Google Test による既存テスト約45件が整備済み", 13)
add_bullet(tf, "", 10)
add_bullet(tf, "3領域 8関数 合計約385行を選定:", 13, bold=True, color=BLACK)
add_bullet(tf, "日時書式・バージョン解析 (format.cpp): 約115行", 13, level=1)
add_bullet(tf, "メールアドレス・文字種判定 (CWordParse.cpp): 約180行", 13, level=1)
add_bullet(tf, "全角半角英数変換 (convert_util.cpp): 約90行", 13, level=1)
add_bullet(tf, "", 10)
add_bullet(tf, "GUI依存・ファイルI/O依存等の11関数は除外", 13, color=GRAY)

# Right: Environment
add_text_box(slide, 17.5, 2.5, 15, 1, "実行環境", 16, bold=True, color=BLACK)
tf = add_text_box(slide, 17.5, 3.5, 15, 4, "", 13, color=GRAY)
add_bullet(tf, "macOS (Apple clang 16.0)", 13, color=BLACK)
add_bullet(tf, "Google Test 1.17 (Homebrew)", 13, color=BLACK)
add_bullet(tf, "Windows互換ヘッダで型定義を代替", 13)
add_bullet(tf, "純粋関数のため移植は型定義のみ", 13)

add_text_box(slide, 17.5, 8, 15, 1, "比較対象", 16, bold=True, color=BLACK)
tf = add_text_box(slide, 17.5, 9, 15, 3, "", 13)
add_bullet(tf, "(A) 既存の手動テスト: 約45件", 13, color=BLACK)
add_bullet(tf, "(C) 提案手法: TRM + LLM生成テスト", 13, bold=True, color=ACCENT)


# =====================================================================
# Slide 6: 結果サマリ（数字で語る）
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, 1.5, 0.8, 31, 1.5, "4. 実施結果: Key Numbers", 24, bold=True, color=WHITE)

metrics = [
    ("99件", "テスト要求 (TRM)", "5種別で構造化定義\n平均3.9行/要求の密度"),
    ("145件", "生成テストケース", "既存テスト45件の約3.2倍\n要求カバー率 100%"),
    ("91.7%", "初回PASS率", "133/145件が初回PASS\nFAIL 12件は全てLLMの期待値誤り"),
    ("100%", "修正後PASS率", "TRM IDで12件を5分類に即時診断\n対象関数のバグは0件"),
    ("61.6%", "既存テスト未カバー率", "99件中61件が既存テスト未カバー\nDP種別は100%欠落"),
    ("65.7%", "コード非専門者可読率", "99件中65件がコード知識なしで\n内容を理解可能"),
]
for i, (num, label, desc) in enumerate(metrics):
    row = i // 3
    col = i % 3
    x = 1.5 + col * 10.5
    y = 3 + row * 7
    add_rect(slide, x, y, 9.5, 2, ACCENT, f"{num}", 28)
    add_text_box(slide, x + 0.3, y + 2.2, 9, 0.8, label, 14, bold=True, color=WHITE)
    add_text_box(slide, x + 0.3, y + 3.2, 9, 2.5, desc, 11, color=GRAY)


# =====================================================================
# Slide 7: テスト要求の種別内訳
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.5, 31, 1.5, "TRM生成結果: 5種別 99件のテスト要求", 22, bold=True, color=ACCENT)
add_image(slide, os.path.join(FIG, "fig2-requirements-breakdown.png"), 1, 2, 15)
add_image(slide, os.path.join(FIG, "fig3-test-comparison.png"), 16.5, 2, 16)

add_text_box(slide, 1.5, 13, 30, 3,
    "BR(分岐網羅)が55.6%と最多 — ソースコードの分岐構造から直接導出\n"
    "ER/DPが少ないのは対象が純粋関数のため（副作用関数では増加が予測される）\n"
    "全領域で既存テストの3.1〜4.1倍のテストケースを体系的に生成",
    13, color=GRAY)


# =====================================================================
# Slide 8: コンパイル・実行結果
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.5, 31, 1.5, "コンパイル・実行結果と障害診断", 22, bold=True, color=ACCENT)
add_image(slide, os.path.join(FIG, "fig4-execution-results.png"), 1, 2, 16)

# FAIL analysis
add_text_box(slide, 18, 2, 14.5, 1, "FAIL 12件の原因分類 (TRM IDで即時診断)", 14, bold=True, color=ACCENT2)

fail_data = [
    ("コンポーネント分割の誤解", "5件", "BR-11〜15"),
    ("数字グルーピングの誤解", "1件", "EC内部"),
    ("文字数カウントミス", "2件", "BR-34, EC"),
    ("API呼び出しバグ", "1件", "EC-16"),
    ("マッピング戻り値の誤解", "3件", "BR-39〜41"),
]
for i, (cause, count, trm) in enumerate(fail_data):
    y = 3.5 + i * 1.8
    add_text_box(slide, 18, y, 8, 0.8, cause, 11, bold=True, color=BLACK)
    add_text_box(slide, 26.5, y, 2, 0.8, count, 11, bold=True, color=ACCENT2)
    add_text_box(slide, 28.5, y, 4, 0.8, trm, 10, color=GRAY)

add_rect(slide, 18, 13, 14.5, 3, RGBColor(0xFD,0xF0,0xF0), "", 12)
add_text_box(slide, 18.5, 13.2, 13.5, 0.7, "P3の解決", 14, bold=True, color=ACCENT2)
tf = add_text_box(slide, 18.5, 14.2, 13.5, 1.5, "", 12, color=BLACK)
add_bullet(tf, "TRMなし: 145件のC++コードを逐一突合 → 長時間", 12, color=GRAY)
add_bullet(tf, "TRMあり: TRM IDから日本語説明+ソース参照で即時診断", 12, bold=True, color=BLACK)


# =====================================================================
# Slide 9: 品質ギャップ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.5, 31, 1.5, "既存テストの品質ギャップ分析 (P2の解決)", 22, bold=True, color=ACCENT)
add_image(slide, os.path.join(FIG, "fig5-quality-gap.png"), 1, 2.5, 18)

# Key findings
add_text_box(slide, 20, 2.5, 12.5, 1, "発見された品質ギャップ", 16, bold=True, color=ACCENT2)
tf = add_text_box(slide, 20, 4, 12.5, 7, "", 13, color=BLACK)
add_bullet(tf, "全要求の61.6%が既存テスト未カバー", 13, bold=True, color=ACCENT2)
add_bullet(tf, "", 8)
add_bullet(tf, "DP (依存切替): 100%未カバー", 13, bold=True)
add_bullet(tf, "  関数間の差異検証・往復変換検証が完全欠落", 11, color=GRAY, level=1)
add_bullet(tf, "", 8)
add_bullet(tf, "BV (境界値): 72.7%未カバー", 13, bold=True)
add_bullet(tf, "  境界値の体系的検証が不足", 11, color=GRAY, level=1)
add_bullet(tf, "", 8)
add_bullet(tf, "EC (同値クラス): 70.4%未カバー", 13, bold=True)
add_bullet(tf, "  入力パターンの網羅的検証が不足", 11, color=GRAY, level=1)
add_bullet(tf, "", 8)
add_bullet(tf, "BR (分岐網羅): 52.7%未カバー", 13)
add_bullet(tf, "  主要パスは比較的カバーされている", 11, color=GRAY, level=1)

add_rect(slide, 1.5, 13.5, 30.5, 2.5, RGBColor(0xE8,0xF0,0xF8), "", 12)
add_text_box(slide, 2, 13.7, 29, 2,
    "TRMがあるからこそ「何の観点が欠けているか」を構造的に把握できる。\n"
    "TRMなしでは「テストが少ない」という漠然とした認識にとどまる。",
    14, bold=True, color=ACCENT)


# =====================================================================
# Slide 10: 結論
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.5, 31, 1.5, "5. 結論: 3つの問題への解決", 24, bold=True, color=ACCENT)

solutions = [
    ("P1 解決", "テスト設計根拠の明示",
     "99件のTRM (5種別) を構造化定義\n"
     "各テストケースにTRM IDを付与\n"
     "テスト設計の妥当性を第三者が検証可能に\n"
     "65.7%がコード非専門者にも可読"),
    ("P2 解決", "網羅性の事前管理",
     "テスト要求カバー率100%を事前に管理\n"
     "既存テストの未カバー61件を種別ごとに特定\n"
     "DP種別100%欠落を発見\n"
     "テスト生成「前」に品質ギャップを把握"),
    ("P3 解決", "テスト失敗の構造的診断",
     "FAIL 12件をTRM IDで5分類に即時特定\n"
     "コード突合不要で原因の所在を把握\n"
     "修正後 145/145件 100% PASS\n"
     "対象関数の実装バグ: 0件"),
]

for i, (pid, title, desc) in enumerate(solutions):
    x = 1.5 + i * 10.5
    add_rect(slide, x, 2.5, 9.5, 1.5, ACCENT, f"{pid}\n{title}", 14)
    add_text_box(slide, x + 0.3, 4.5, 9, 5, desc, 12, color=BLACK)

# Future work
add_text_box(slide, 1.5, 10.5, 31, 1, "今後の課題", 16, bold=True, color=BLACK)
tf = add_text_box(slide, 1.5, 11.5, 31, 4, "", 12, color=GRAY)
add_bullet(tf, "(1) 副作用を持つ関数への拡張（モック・スタブ戦略）", 12)
add_bullet(tf, "(2) コードカバレッジ (C0/C1) との相関分析", 12)
add_bullet(tf, "(3) 非開発者による可読性の被験者評価実験", 12)
add_bullet(tf, "(4) 他言語・他プロジェクトでの再現実験", 12)


# =====================================================================
# Slide 11: 先行研究との位置づけ
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 1.5, 0.5, 31, 1.5, "付録: 先行研究との位置づけ", 22, bold=True, color=ACCENT)

add_text_box(slide, 1.5, 2.5, 31, 1,
    "LLMテスト生成研究の全てが「テスト生成精度の向上」に集中 → テスト設計情報の外部化は未開拓",
    14, bold=True, color=ACCENT2)

# Table-like comparison
headers = ["手法", "年", "主眼", "テスト設計の可視化"]
rows = [
    ["CodaMOSA [1]", "ICSE 2023", "SBST+LLMハイブリッド", "扱わない"],
    ["ChatUniTest [2]", "FSE 2024", "生成-検証-修復サイクル", "扱わない"],
    ["TestPilot [3]", "IEEE TSE 2024", "ゼロショット+再プロンプト", "扱わない"],
    ["CoverUp [4]", "2024", "カバレッジ誘導型生成", "扱わない"],
    ["TELPA [6]", "ACM TOSEM 2024", "静的解析+困難分岐特化", "扱わない"],
    ["本研究", "SQiP 2026", "TRMによる構造化+生成", "主要な貢献"],
]

for j, h in enumerate(headers):
    x = 1.5 + j * (7.5 if j < 3 else 6)
    w = 7 if j < 3 else 5.5
    if j == 0: x, w = 1.5, 5.5
    elif j == 1: x, w = 7.5, 4
    elif j == 2: x, w = 12, 10
    elif j == 3: x, w = 22.5, 10
    add_text_box(slide, x, 4, w, 0.8, h, 11, bold=True, color=WHITE)
    add_rect(slide, x, 3.8, w, 0.9, ACCENT, h, 11)

for i, row in enumerate(rows):
    y = 5.2 + i * 1.5
    is_ours = i == len(rows) - 1
    for j, cell in enumerate(row):
        if j == 0: x, w = 1.5, 5.5
        elif j == 1: x, w = 7.5, 4
        elif j == 2: x, w = 12, 10
        elif j == 3: x, w = 22.5, 10
        c = ACCENT if (is_ours and j == 3) else (ACCENT2 if is_ours else BLACK)
        add_text_box(slide, x, y, w, 0.8, cell, 11, bold=is_ours, color=c)


# =====================================================================
# Save
# =====================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(OUT):,} bytes")
